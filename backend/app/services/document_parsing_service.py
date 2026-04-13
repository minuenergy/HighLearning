from __future__ import annotations

import json
import os
import subprocess
import tempfile
from functools import lru_cache
from pathlib import Path
from typing import Any

import pymupdf

from app.config import settings
from app.utils.ocr_runtime import default_document_ocr_backend, is_macos, resolve_paddle_device


class DocumentParsingError(RuntimeError):
    pass


def _page_row(page_number: int, text: str, *, label: str | None = None) -> dict[str, Any]:
    normalized = text.strip()
    return {
        "page_number": page_number,
        "page_label": label or f"{page_number}p",
        "text": normalized,
    }


def _merge_page_rows(rows: list[dict[str, Any]]) -> str:
    return "\n\n".join(row["text"] for row in rows if row.get("text"))


def _extract_native_pdf_text(file_bytes: bytes) -> str:
    document = pymupdf.open(stream=file_bytes, filetype="pdf")
    return "\n".join(page.get_text() for page in document)


def _extract_native_pdf_pages(file_bytes: bytes) -> list[dict[str, Any]]:
    document = pymupdf.open(stream=file_bytes, filetype="pdf")
    rows: list[dict[str, Any]] = []
    for page_number, page in enumerate(document, start=1):
        rows.append(_page_row(page_number, page.get_text()))
    return rows


def _should_use_paddleocr(native_text: str | int, parser_mode: str) -> bool:
    if parser_mode == "paddleocr":
        return True
    if parser_mode == "pymupdf":
        return False
    if parser_mode == "visionocr":
        return True
    if isinstance(native_text, int):
        native_chars = native_text
    else:
        native_chars = len(native_text.strip())
    return native_chars < settings.document_parser_native_min_chars


def _result_field(result: Any, field_name: str) -> Any:
    if hasattr(result, field_name):
        return getattr(result, field_name)
    if isinstance(result, dict):
        return result.get(field_name)
    return None


def _batched(items: list[Path], batch_size: int) -> list[list[Path]]:
    safe_batch_size = max(1, batch_size)
    return [items[index : index + safe_batch_size] for index in range(0, len(items), safe_batch_size)]


@lru_cache(maxsize=1)
def _get_paddleocr_pipeline():
    try:
        from paddleocr import PPStructureV3
    except ImportError as error:
        raise DocumentParsingError(
            "PaddleOCR가 설치되어 있지 않습니다. PaddleOCR 공식 문서에 따라 "
            "PaddlePaddle과 paddleocr[all]을 설치한 뒤 다시 시도해주세요."
        ) from error

    return PPStructureV3(
        lang=settings.paddleocr_lang,
        use_doc_orientation_classify=settings.paddleocr_use_doc_orientation_classify,
        use_doc_unwarping=settings.paddleocr_use_doc_unwarping,
        use_textline_orientation=settings.paddleocr_use_textline_orientation,
        device=resolve_paddle_device(settings.paddleocr_device),
    )


def _vision_ocr_paths() -> tuple[Path, Path]:
    backend_root = Path(__file__).resolve().parents[2]
    return (
        backend_root / "scripts" / "vision_ocr.swift",
        backend_root / ".vision-ocr" / "vision_ocr",
    )


def _ensure_vision_ocr_binary() -> Path:
    source, binary = _vision_ocr_paths()
    if not source.exists():
        raise DocumentParsingError(f"Vision OCR 소스 파일을 찾지 못했습니다: {source}")
    if not is_macos():
        raise DocumentParsingError("Vision OCR는 macOS에서만 지원됩니다. Windows에서는 PaddleOCR를 사용해주세요.")
    if binary.exists() and binary.stat().st_mtime >= source.stat().st_mtime:
        return binary

    binary.parent.mkdir(parents=True, exist_ok=True)
    module_cache_dir = Path(__file__).resolve().parents[2] / ".swift-module-cache"
    module_cache_dir.mkdir(parents=True, exist_ok=True)
    compile_env = os.environ.copy()
    compile_env.setdefault("SWIFT_MODULECACHE_PATH", str(module_cache_dir))
    compile_env.setdefault("CLANG_MODULE_CACHE_PATH", str(module_cache_dir))
    subprocess.run(
        ["swiftc", "-O", str(source), "-o", str(binary)],
        check=True,
        capture_output=True,
        text=True,
        env=compile_env,
    )
    return binary


def _run_vision_ocr_batch(image_paths: list[Path]) -> dict[str, str]:
    binary = _ensure_vision_ocr_binary()
    languages = [item.strip() for item in settings.vision_ocr_languages.split(",") if item.strip()]
    if not languages:
        languages = ["ko-KR", "en-US"]

    with tempfile.NamedTemporaryFile("w", encoding="utf-8", delete=False) as handle:
        for path in image_paths:
            handle.write(str(path.resolve()) + "\n")
        input_list_path = Path(handle.name)

    try:
        command = [
            str(binary),
            "--input-list",
            str(input_list_path),
            "--langs",
            ",".join(languages),
        ]
        if settings.vision_ocr_fast:
            command.append("--fast")
        process = subprocess.run(
            command,
            check=True,
            capture_output=True,
            text=True,
        )
    except subprocess.CalledProcessError as error:
        raise DocumentParsingError(error.stderr.strip() or "Vision OCR 실행에 실패했습니다.") from error
    finally:
        input_list_path.unlink(missing_ok=True)

    results: dict[str, str] = {}
    for line in process.stdout.splitlines():
        stripped = line.strip()
        if not stripped:
            continue
        payload = json.loads(stripped)
        if payload.get("error"):
            raise DocumentParsingError(f"Vision OCR failed for {payload.get('path')}: {payload['error']}")
        results[payload["path"]] = payload.get("text", "")
    return results


def _run_vision_ocr_batched(image_paths: list[Path]) -> dict[str, str]:
    results: dict[str, str] = {}
    for batch in _batched(image_paths, settings.document_parser_ocr_batch_size):
        results.update(_run_vision_ocr_batch(batch))
    return results


def _extract_pdf_with_visionocr(file_bytes: bytes) -> str:
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        document = pymupdf.open(stream=file_bytes, filetype="pdf")
        image_paths: list[Path] = []

        for page_index, page in enumerate(document, start=1):
            pixmap = page.get_pixmap(dpi=settings.vision_ocr_pdf_dpi, alpha=False)
            image_path = temp_path / f"page{page_index:05d}.png"
            pixmap.save(str(image_path))
            image_paths.append(image_path)

        texts = _run_vision_ocr_batched(image_paths)
        merged = "\n\n".join(
            text.strip()
            for path in image_paths
            if (text := texts.get(str(path.resolve()), "")).strip()
        )
        if merged.strip():
            return merged

    raise DocumentParsingError("Vision OCR 파싱 결과에서 사용할 텍스트를 찾지 못했습니다.")


def _extract_pdf_pages_with_visionocr(file_bytes: bytes) -> list[dict[str, Any]]:
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        document = pymupdf.open(stream=file_bytes, filetype="pdf")
        image_paths: list[Path] = []

        for page_index, page in enumerate(document, start=1):
            pixmap = page.get_pixmap(dpi=settings.vision_ocr_pdf_dpi, alpha=False)
            image_path = temp_path / f"page{page_index:05d}.png"
            pixmap.save(str(image_path))
            image_paths.append(image_path)

        texts = _run_vision_ocr_batched(image_paths)
        rows = [
            _page_row(index, texts.get(str(path.resolve()), ""))
            for index, path in enumerate(image_paths, start=1)
            if texts.get(str(path.resolve()), "").strip()
        ]
        if rows:
            return rows

    raise DocumentParsingError("Vision OCR 페이지 파싱 결과에서 사용할 텍스트를 찾지 못했습니다.")


def _extract_pdf_with_paddleocr(file_bytes: bytes) -> str:
    pipeline = _get_paddleocr_pipeline()

    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=True) as temp_file:
        temp_file.write(file_bytes)
        temp_file.flush()

        output = pipeline.predict(input=temp_file.name)
        markdown_pages: list[Any] = []
        text_fragments: list[str] = []

        for result in output:
            markdown = _result_field(result, "markdown")
            if markdown:
                markdown_pages.append(markdown)
                if isinstance(markdown, dict):
                    markdown_text = markdown.get("text") or markdown.get("markdown_text")
                    if markdown_text:
                        text_fragments.append(markdown_text)
                elif isinstance(markdown, str):
                    text_fragments.append(markdown)

            parsing_blocks = _result_field(result, "parsing_res_list") or []
            for block in parsing_blocks:
                block_content = block.get("block_content")
                if block_content:
                    text_fragments.append(block_content)

        if markdown_pages and hasattr(pipeline, "concatenate_markdown_pages"):
            merged_markdown = pipeline.concatenate_markdown_pages(markdown_pages)
            if isinstance(merged_markdown, str) and merged_markdown.strip():
                return merged_markdown

        fallback_text = "\n\n".join(fragment.strip() for fragment in text_fragments if fragment and fragment.strip())
        if fallback_text.strip():
            return fallback_text

    raise DocumentParsingError("PaddleOCR 파싱 결과에서 사용할 텍스트를 찾지 못했습니다.")


def _extract_pdf_pages_with_paddleocr(file_bytes: bytes) -> list[dict[str, Any]]:
    pipeline = _get_paddleocr_pipeline()

    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=True) as temp_file:
        temp_file.write(file_bytes)
        temp_file.flush()

        output = pipeline.predict(input=temp_file.name)
        page_rows: list[dict[str, Any]] = []

        for page_number, result in enumerate(output, start=1):
            text_fragments: list[str] = []
            markdown = _result_field(result, "markdown")
            if markdown:
                if isinstance(markdown, dict):
                    markdown_text = markdown.get("text") or markdown.get("markdown_text")
                    if markdown_text:
                        text_fragments.append(markdown_text)
                elif isinstance(markdown, str):
                    text_fragments.append(markdown)

            parsing_blocks = _result_field(result, "parsing_res_list") or []
            for block in parsing_blocks:
                block_content = block.get("block_content")
                if block_content:
                    text_fragments.append(block_content)

            page_text = "\n\n".join(fragment.strip() for fragment in text_fragments if fragment and fragment.strip()).strip()
            if page_text:
                page_rows.append(_page_row(page_number, page_text))

        if page_rows:
            return page_rows

    raise DocumentParsingError("PaddleOCR 페이지 파싱 결과에서 사용할 텍스트를 찾지 못했습니다.")


def extract_pdf_text(file_bytes: bytes, parser_mode: str | None = None) -> tuple[str, str]:
    mode = normalize_parser_mode(parser_mode or settings.document_parser)
    native_text = _extract_native_pdf_text(file_bytes)

    if not _should_use_paddleocr(native_text, mode):
        return native_text, "pymupdf"

    fallback_backend = default_document_ocr_backend() if mode == "auto" else mode

    if fallback_backend == "visionocr":
        try:
            vision_text = _extract_pdf_with_visionocr(file_bytes)
            if vision_text.strip():
                return vision_text, "visionocr"
        except DocumentParsingError:
            if mode == "visionocr":
                raise

    try:
        paddle_text = _extract_pdf_with_paddleocr(file_bytes)
        if paddle_text.strip():
            return paddle_text, "paddleocr"
    except DocumentParsingError:
        if mode == "paddleocr":
            raise

    return native_text, "pymupdf"


def extract_pdf_pages(file_bytes: bytes, parser_mode: str | None = None) -> tuple[list[dict[str, Any]], str]:
    mode = normalize_parser_mode(parser_mode or settings.document_parser)
    native_rows = _extract_native_pdf_pages(file_bytes)
    native_char_count = sum(len(str(row.get("text") or "").strip()) for row in native_rows)

    if not _should_use_paddleocr(native_char_count, mode):
        return native_rows, "pymupdf"

    fallback_backend = default_document_ocr_backend() if mode == "auto" else mode

    if fallback_backend == "visionocr":
        try:
            vision_rows = _extract_pdf_pages_with_visionocr(file_bytes)
            if vision_rows:
                return vision_rows, "visionocr"
        except DocumentParsingError:
            if mode == "visionocr":
                raise

    try:
        paddle_rows = _extract_pdf_pages_with_paddleocr(file_bytes)
        if paddle_rows:
            return paddle_rows, "paddleocr"
    except DocumentParsingError:
        if mode == "paddleocr":
            raise

    return native_rows, "pymupdf"


def extract_pptx_pages(file_bytes: bytes) -> tuple[list[dict[str, Any]], str]:
    import io
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(file_bytes))
    rows: list[dict[str, Any]] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        slide_text = "\n".join(texts).strip()
        if slide_text:
            rows.append(_page_row(slide_number, slide_text, label=f"{slide_number}slide"))
    return rows, "python-pptx"


def extract_pptx_text(file_bytes: bytes) -> tuple[str, str]:
    rows, parser_used = extract_pptx_pages(file_bytes)
    return _merge_page_rows(rows), parser_used


def _extract_native_pdf_text_from_path(file_path: str | Path) -> str:
    document = pymupdf.open(str(file_path))
    return "\n".join(page.get_text() for page in document)


def _extract_native_pdf_pages_from_path(file_path: str | Path) -> list[dict[str, Any]]:
    document = pymupdf.open(str(file_path))
    rows: list[dict[str, Any]] = []
    for page_number, page in enumerate(document, start=1):
        rows.append(_page_row(page_number, page.get_text()))
    return rows


def _extract_pdf_with_visionocr_from_path(file_path: str | Path) -> str:
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        document = pymupdf.open(str(file_path))
        image_paths: list[Path] = []

        for page_index, page in enumerate(document, start=1):
            pixmap = page.get_pixmap(dpi=settings.vision_ocr_pdf_dpi, alpha=False)
            image_path = temp_path / f"page{page_index:05d}.png"
            pixmap.save(str(image_path))
            image_paths.append(image_path)

        texts = _run_vision_ocr_batched(image_paths)
        merged = "\n\n".join(
            text.strip()
            for path in image_paths
            if (text := texts.get(str(path.resolve()), "")).strip()
        )
        if merged.strip():
            return merged

    raise DocumentParsingError("Vision OCR 파싱 결과에서 사용할 텍스트를 찾지 못했습니다.")


def _extract_pdf_pages_with_visionocr_from_path(file_path: str | Path) -> list[dict[str, Any]]:
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        document = pymupdf.open(str(file_path))
        image_paths: list[Path] = []

        for page_index, page in enumerate(document, start=1):
            pixmap = page.get_pixmap(dpi=settings.vision_ocr_pdf_dpi, alpha=False)
            image_path = temp_path / f"page{page_index:05d}.png"
            pixmap.save(str(image_path))
            image_paths.append(image_path)

        texts = _run_vision_ocr_batched(image_paths)
        rows = [
            _page_row(index, texts.get(str(path.resolve()), ""))
            for index, path in enumerate(image_paths, start=1)
            if texts.get(str(path.resolve()), "").strip()
        ]
        if rows:
            return rows

    raise DocumentParsingError("Vision OCR 페이지 파싱 결과에서 사용할 텍스트를 찾지 못했습니다.")


def _extract_pdf_with_paddleocr_from_path(file_path: str | Path) -> str:
    pipeline = _get_paddleocr_pipeline()
    output = pipeline.predict(input=str(file_path))
    markdown_pages: list[Any] = []
    text_fragments: list[str] = []

    for result in output:
        markdown = _result_field(result, "markdown")
        if markdown:
            markdown_pages.append(markdown)
            if isinstance(markdown, dict):
                markdown_text = markdown.get("text") or markdown.get("markdown_text")
                if markdown_text:
                    text_fragments.append(markdown_text)
            elif isinstance(markdown, str):
                text_fragments.append(markdown)

        parsing_blocks = _result_field(result, "parsing_res_list") or []
        for block in parsing_blocks:
            block_content = block.get("block_content")
            if block_content:
                text_fragments.append(block_content)

    if markdown_pages and hasattr(pipeline, "concatenate_markdown_pages"):
        merged_markdown = pipeline.concatenate_markdown_pages(markdown_pages)
        if isinstance(merged_markdown, str) and merged_markdown.strip():
            return merged_markdown

    fallback_text = "\n\n".join(fragment.strip() for fragment in text_fragments if fragment and fragment.strip())
    if fallback_text.strip():
        return fallback_text

    raise DocumentParsingError("PaddleOCR 파싱 결과에서 사용할 텍스트를 찾지 못했습니다.")


def _extract_pdf_pages_with_paddleocr_from_path(file_path: str | Path) -> list[dict[str, Any]]:
    pipeline = _get_paddleocr_pipeline()
    output = pipeline.predict(input=str(file_path))
    page_rows: list[dict[str, Any]] = []

    for page_number, result in enumerate(output, start=1):
        text_fragments: list[str] = []
        markdown = _result_field(result, "markdown")
        if markdown:
            if isinstance(markdown, dict):
                markdown_text = markdown.get("text") or markdown.get("markdown_text")
                if markdown_text:
                    text_fragments.append(markdown_text)
            elif isinstance(markdown, str):
                text_fragments.append(markdown)

        parsing_blocks = _result_field(result, "parsing_res_list") or []
        for block in parsing_blocks:
            block_content = block.get("block_content")
            if block_content:
                text_fragments.append(block_content)

        page_text = "\n\n".join(fragment.strip() for fragment in text_fragments if fragment and fragment.strip()).strip()
        if page_text:
            page_rows.append(_page_row(page_number, page_text))

    if page_rows:
        return page_rows

    raise DocumentParsingError("PaddleOCR 페이지 파싱 결과에서 사용할 텍스트를 찾지 못했습니다.")


def extract_pdf_text_from_path(file_path: str | Path, parser_mode: str | None = None) -> tuple[str, str]:
    mode = normalize_parser_mode(parser_mode or settings.document_parser)
    native_text = _extract_native_pdf_text_from_path(file_path)

    if not _should_use_paddleocr(native_text, mode):
        return native_text, "pymupdf"

    fallback_backend = default_document_ocr_backend() if mode == "auto" else mode

    if fallback_backend == "visionocr":
        try:
            vision_text = _extract_pdf_with_visionocr_from_path(file_path)
            if vision_text.strip():
                return vision_text, "visionocr"
        except DocumentParsingError:
            if mode == "visionocr":
                raise

    try:
        paddle_text = _extract_pdf_with_paddleocr_from_path(file_path)
        if paddle_text.strip():
            return paddle_text, "paddleocr"
    except DocumentParsingError:
        if mode == "paddleocr":
            raise

    return native_text, "pymupdf"


def extract_pdf_pages_from_path(file_path: str | Path, parser_mode: str | None = None) -> tuple[list[dict[str, Any]], str]:
    mode = normalize_parser_mode(parser_mode or settings.document_parser)
    native_rows = _extract_native_pdf_pages_from_path(file_path)
    native_char_count = sum(len(str(row.get("text") or "").strip()) for row in native_rows)

    if not _should_use_paddleocr(native_char_count, mode):
        return native_rows, "pymupdf"

    fallback_backend = default_document_ocr_backend() if mode == "auto" else mode

    if fallback_backend == "visionocr":
        try:
            vision_rows = _extract_pdf_pages_with_visionocr_from_path(file_path)
            if vision_rows:
                return vision_rows, "visionocr"
        except DocumentParsingError:
            if mode == "visionocr":
                raise

    try:
        paddle_rows = _extract_pdf_pages_with_paddleocr_from_path(file_path)
        if paddle_rows:
            return paddle_rows, "paddleocr"
    except DocumentParsingError:
        if mode == "paddleocr":
            raise

    return native_rows, "pymupdf"


def extract_pptx_pages_from_path(file_path: str | Path) -> tuple[list[dict[str, Any]], str]:
    from pptx import Presentation

    presentation = Presentation(str(file_path))
    rows: list[dict[str, Any]] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
        slide_text = "\n".join(texts).strip()
        if slide_text:
            rows.append(_page_row(slide_number, slide_text, label=f"{slide_number}slide"))
    return rows, "python-pptx"


def normalize_parser_mode(parser_mode: str) -> str:
    normalized = parser_mode.lower()
    if normalized not in {"auto", "pymupdf", "paddleocr", "visionocr"}:
        raise DocumentParsingError(
            f"지원하지 않는 parser_mode입니다: {parser_mode}. 가능한 값은 auto, pymupdf, paddleocr, visionocr 입니다."
        )
    return normalized


def parse_document(file_name: str, file_bytes: bytes, parser_mode: str | None = None) -> tuple[str, str]:
    page_rows, parser_used = parse_document_pages(file_name, file_bytes, parser_mode)
    return _merge_page_rows(page_rows), parser_used


def parse_document_pages(
    file_name: str,
    file_bytes: bytes,
    parser_mode: str | None = None,
) -> tuple[list[dict[str, Any]], str]:
    suffix = Path(file_name).suffix.lower()

    if suffix == ".pptx":
        return extract_pptx_pages(file_bytes)
    if suffix == ".pdf":
        return extract_pdf_pages(file_bytes, parser_mode)

    raise DocumentParsingError(f"지원하지 않는 파일 형식입니다: {suffix or file_name}")


def parse_document_pages_from_path(
    file_name: str,
    file_path: str | Path,
    parser_mode: str | None = None,
) -> tuple[list[dict[str, Any]], str]:
    suffix = Path(file_name).suffix.lower()

    if suffix == ".pptx":
        return extract_pptx_pages_from_path(file_path)
    if suffix == ".pdf":
        return extract_pdf_pages_from_path(file_path, parser_mode)

    raise DocumentParsingError(f"지원하지 않는 파일 형식입니다: {suffix or file_name}")
